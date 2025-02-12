import logging
import pandas as pd
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
import os

# Configure logging
logging.basicConfig(level=logging.DEBUG)
logger = logging.getLogger(__name__)


def validate_excel_data(df):
    """
    Validate the Excel data matches the expected format from VBA code
    """
    logger.debug("Validating Excel data structure")

    # Check if required columns exist (using column indices from VBA code)
    if len(df.columns) < 9:
        error_msg = "Excel file must have at least 9 columns"
        logger.error(error_msg)
        raise ValueError(error_msg)

    # Rename columns to match our processing logic
    df = df.rename(
        columns={
            df.columns[5]: 'Description',  # Column 6 in VBA
            df.columns[6]: 'Store Price',  # Column 7 in VBA
            df.columns[7]: 'Offer Price',  # Column 8 in VBA
            df.columns[8]: 'Discount'  # Column 9 in VBA
        })

    # Clean and validate data types
    try:
        logger.debug("Converting data types")
        # Clean price columns by removing 'Leke' and converting to numeric
        df['Store Price'] = df['Store Price'].astype(str).str.replace(
            ' Leke', '', regex=False)
        df['Store Price'] = pd.to_numeric(df['Store Price'])

        df['Offer Price'] = df['Offer Price'].astype(str).str.replace(
            ' Leke', '', regex=False)
        df['Offer Price'] = pd.to_numeric(df['Offer Price'])

        # Clean discount column and convert to proper percentage
        df['Discount'] = df['Discount'].astype(str).str.rstrip('%')
        df['Discount'] = pd.to_numeric(df['Discount'])
        # Multiply by 100 if the value is less than 1 (converting decimal to percentage)
        df.loc[df['Discount'] < 1, 'Discount'] *= 100

    except Exception as e:
        error_msg = f"Error converting data types: {str(e)}"
        logger.error(error_msg)
        raise ValueError(error_msg)

    return df


def format_price(price):
    """
    Format price with thousand separators (matches VBA's Format(value, "#,##0"))
    """
    return f"{price:,.0f}"


def update_shape_text(shape, text, font_size=None, font_color=None):
    """
    Update PowerPoint shape text with formatting using python-pptx
    """
    try:
        if not hasattr(shape, 'text_frame'):
            logger.warning(f"Shape {shape.name} does not have a text frame")
            return

        text_frame = shape.text_frame
        text_frame.clear()  # Clear existing text
        p = text_frame.paragraphs[0]
        p.text = str(text)

        if font_size:
            p.font.size = Pt(font_size)
        if font_color:
            p.font.color.rgb = RGBColor(*font_color)

        # Make text bold
        p.font.bold = True

    except Exception as e:
        logger.error(f"Error updating shape text: {str(e)}")
        raise


def find_shape_by_name(slide, name_pattern):
    """
    Find a shape by its name pattern in a slide
    """
    for shape in slide.shapes:
        if hasattr(shape,
                   'name') and name_pattern.lower() in shape.name.lower():
            return shape
    return None


def process_powerpoint_slide(slide, row_data, sticker_index):
    """
    Process a single PowerPoint slide with the given row data
    """
    try:
        logger.debug(f"Processing slide for sticker index {sticker_index}")

        # Format the values (matching VBA formatting)
        offer_price = format_price(row_data['Offer Price'])
        store_price = format_price(row_data['Store Price'])
        discount = f"{int(row_data['Discount'])}%"  # Convert to integer to remove decimal places

        # Update shapes for this sticker (matching VBA naming)
        shapes_updated = False

        # Product description
        shape = find_shape_by_name(slide,
                                   f"pershkrimi i produktit{sticker_index}")
        if shape:
            update_shape_text(shape, row_data['Description'], font_size=28)
            shapes_updated = True

        # Offer price
        shape = find_shape_by_name(slide, f"cmim oferte{sticker_index}")
        if shape:
            update_shape_text(shape,
                              offer_price,
                              font_size=48,
                              font_color=(0, 0, 0))  # Changed from 72 to 35
            shapes_updated = True

        # Store price
        shape = find_shape_by_name(slide, f"cmim dyqani{sticker_index}")
        if shape:
            update_shape_text(shape,
                              store_price,
                              font_size=28,
                              font_color=(0, 0, 0))  # Changed from 45 to 12
            shapes_updated = True

        # Discount percentage
        shape = find_shape_by_name(slide, f"ulja ne %{sticker_index}")
        if shape:
            update_shape_text(shape,
                              discount,
                              font_size=22,
                              font_color=(255, 255, 255))
            shapes_updated = True

        if not shapes_updated:
            logger.warning(
                f"No shapes were updated for sticker index {sticker_index}")

    except Exception as e:
        logger.error(
            f"Error processing slide for sticker {sticker_index}: {str(e)}")
        raise


def transfer_data_to_powerpoint(excel_path, ppt_path, output_path):
    """
    Main function to transfer data from Excel to PowerPoint
    """
    try:
        logger.info("Starting data transfer process")
        logger.debug(f"Reading Excel file from: {excel_path}")

        # Read Excel file (start from row 2 as in VBA code)
        df = pd.read_excel(excel_path, header=0)
        df = validate_excel_data(df)

        logger.debug(f"Opening PowerPoint file from: {ppt_path}")
        # Open PowerPoint presentation
        prs = Presentation(ppt_path)

        # Process each row in Excel (starting from index 0 which is row 2 in Excel)
        for index, row in df.iterrows():
            slide_index = index // 3  # 3 stickers per slide
            sticker_index = (index % 3) + 1  # 1-based sticker index as in VBA

            logger.debug(
                f"Processing row {index + 1} on slide {slide_index + 1}, sticker {sticker_index}"
            )

            # Add new slide if needed
            while slide_index >= len(prs.slides):
                prs.slides.add_slide(prs.slide_layouts[0])

            # Process the slide
            process_powerpoint_slide(prs.slides[slide_index], row,
                                     sticker_index)

        logger.debug(f"Saving modified presentation to: {output_path}")
        # Save the modified presentation
        prs.save(output_path)
        logger.info("PowerPoint file processed successfully")

        return True

    except Exception as e:
        logger.error(f"Error in transfer_data_to_powerpoint: {str(e)}")
        raise


def transfer_stickers_to_powerpoint(excel_path, ppt_path, output_path):
    """
    Transfer sticker data from Excel to PowerPoint using the VBA code logic
    """
    try:
        logger.info("Starting stickers transfer process")
        logger.debug(f"Reading Excel file from: {excel_path}")

        # Read Excel file (start from row 2 as in VBA code)
        df = pd.read_excel(excel_path, header=0)

        # Validate columns exist (using column indices from updated VBA code)
        if len(df.columns) < 5:
            error_msg = "Excel file must have at least 5 columns"
            logger.error(error_msg)
            raise ValueError(error_msg)

        # Rename columns to match VBA code
        df = df.rename(
            columns={
                df.columns[2]: 'Description',  # Column 3 in VBA
                df.columns[3]: 'Store Price',  # Column 4 in VBA
                df.columns[4]: 'Offer Price',  # Column 5 in VBA
            })

        logger.debug(f"Opening PowerPoint file from: {ppt_path}")
        # Open PowerPoint presentation
        prs = Presentation(ppt_path)

        # Process each row in Excel (starting from index 0 which is row 2 in Excel)
        for index, row in df.iterrows():
            slide_index = index // 9  # 9 stickers per slide
            sticker_index = (index % 9) + 1  # 1-based sticker index as in VBA

            logger.debug(
                f"Processing row {index + 1} on slide {slide_index + 1}, sticker {sticker_index}"
            )

            # Add new slide if needed
            while slide_index >= len(prs.slides):
                prs.slides.add_slide(prs.slide_layouts[0])

            # Process the slide with sticker data
            process_sticker_slide(prs.slides[slide_index], row, sticker_index)

        logger.debug(f"Saving modified presentation to: {output_path}")
        # Save the modified presentation
        prs.save(output_path)
        logger.info("PowerPoint stickers file processed successfully")

        return True

    except Exception as e:
        logger.error(f"Error in transfer_stickers_to_powerpoint: {str(e)}")
        raise


def process_sticker_slide(slide, row_data, sticker_index):
    """
    Process a single PowerPoint slide with the given sticker data
    """
    try:
        logger.debug(f"Processing slide for sticker index {sticker_index}")

        # Clean and format the price values (following VBA logic)
        store_price = str(row_data['Store Price']).replace(' Leke', '')
        offer_price = str(row_data['Offer Price']).replace(' Leke', '')

        try:
            store_price = format_price(float(store_price))
        except (ValueError, TypeError):
            store_price = ""

        try:
            offer_price = format_price(float(offer_price))
        except (ValueError, TypeError):
            offer_price = ""

        # Update shapes for this sticker
        shapes_updated = False

        # Product description
        shape = find_shape_by_name(slide,
                                   f"pershkrimi i produktit{sticker_index}")
        if shape:
            update_shape_text(shape,
                              row_data['Description'],
                              font_size=11,
                              font_color=(0, 0, 0))
            shapes_updated = True

        # Store price
        shape = find_shape_by_name(slide, f"cmim dyqani{sticker_index}")
        if shape:
            update_shape_text(shape,
                              store_price,
                              font_size=12,
                              font_color=(0, 0, 0))
            shapes_updated = True

        # Offer price
        shape = find_shape_by_name(slide, f"cmim oferte{sticker_index}")
        if shape:
            update_shape_text(shape,
                              offer_price,
                              font_size=35,
                              font_color=(0, 0, 0))
            shapes_updated = True

        if not shapes_updated:
            logger.warning(
                f"No shapes were updated for sticker index {sticker_index}")

    except Exception as e:
        logger.error(
            f"Error processing slide for sticker {sticker_index}: {str(e)}")
        raise
